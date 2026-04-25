from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()

# Light green RGB (e.g., #E8F5E9)
light_green = RGBColor(232, 245, 233)
dark_green = RGBColor(30, 80, 40)

def apply_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = light_green

slides_data = [
    {
        "title": "Feature Model to Alloy Verification",
        "content": "A Minimal Proof of Concept\n\nSemester Project Presentation",
        "is_title": True
    },
    {
        "title": "Our Objective",
        "content": [
            "Provide a fully automated pipeline from a Feature Model to formal verification.",
            "Keep the implementation minimal, clear, and fully explainable.",
            "Prove that product line constraints can be validated using the Alloy Analyzer."
        ]
    },
    {
        "title": "The Pipeline Architecture",
        "is_flowchart": True
    },
    {
        "title": "A Simple, Flat Grammar (Xtext)",
        "content": [
            "Designed for maximum clarity with only 5 core rules.",
            "Features are defined as 'must' (mandatory) or 'may' (optional).",
            "Constraints use 'requires' and 'excludes' rules between features.",
            "Example: 'must Security' and 'Lights requires WiFi'."
        ]
    },
    {
        "title": "The Code Generator (Xtend)",
        "content": [
            "Automatically triggered when the .fm file is saved in Eclipse.",
            "Assembles an Alloy module and generates 'one sig' for each feature.",
            "Translates mandatory rules into 'X in s'.",
            "Translates constraints into Alloy logic (e.g., 'implies' and 'not')."
        ]
    },
    {
        "title": "Act 1: A Valid Model",
        "content": [
            "We define a 'SmartHome' with Security (must), WiFi (must), and Lights (may).",
            "Constraint: Lights requires WiFi.",
            "Alloy translates this and searches for a valid product.",
            "Result: \"Instance found.\" The model is consistent and satisfiable!"
        ]
    },
    {
        "title": "Act 2: Catching Errors",
        "content": [
            "We introduce a deliberate contradiction: 'Security excludes Security'.",
            "However, Security is also defined as a 'must' (mandatory) feature.",
            "Both cannot be true at the same time.",
            "Result: \"No instance found.\" The analyzer successfully catches the flaw."
        ]
    }
]

for slide_data in slides_data:
    if slide_data.get("is_title"):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        apply_background(slide)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = slide_data["title"]
        subtitle.text = slide_data["content"]
        
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = dark_green
        
    elif slide_data.get("is_flowchart"):
        slide_layout = prs.slide_layouts[5] # Title only
        slide = prs.slides.add_slide(slide_layout)
        apply_background(slide)
        title = slide.shapes.title
        title.text = slide_data["title"]
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = dark_green
        
        # Draw Flowchart
        stages = [
            ("1. DSL (.fm)", "Define the feature model and constraints in plain text."),
            ("2. Xtext Parser", "Parses the text and converts it into an EMF Model."),
            ("3. Xtend Generator", "Translates the EMF Model into Alloy code (.als)."),
            ("4. Alloy Analyzer", "Executes the .als file to verify satisfiability.")
        ]
        
        box_width = Inches(2.7)
        box_height = Inches(0.8)
        left_margin = Inches(0.8)
        start_top = Inches(1.8)
        vertical_spacing = Inches(1.3)
        
        for i, (stage_name, stage_desc) in enumerate(stages):
            top = start_top + (i * vertical_spacing)
            
            # Draw Box
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, left_margin, top, box_width, box_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(190, 230, 190)
            shape.line.color.rgb = dark_green
            shape.line.width = Pt(1.5)
            
            tf = shape.text_frame
            tf.text = stage_name
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = dark_green
            tf.paragraphs[0].font.size = Pt(22)
            
            # Draw description text box next to it
            desc_left = left_margin + box_width + Inches(0.4)
            desc_width = Inches(5.8)
            desc_box = slide.shapes.add_textbox(desc_left, top, desc_width, box_height)
            desc_tf = desc_box.text_frame
            desc_tf.word_wrap = True
            p = desc_tf.paragraphs[0]
            p.text = stage_desc
            p.font.size = Pt(22)
            p.font.color.rgb = RGBColor(50, 50, 50)
            desc_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # Draw down arrow (except for the last one)
            if i < len(stages) - 1:
                arrow_top = top + box_height + Inches(0.05)
                arrow_left = left_margin + box_width/2 - Inches(0.15)
                arrow_width = Inches(0.3)
                arrow_height = vertical_spacing - box_height - Inches(0.1)
                
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.DOWN_ARROW, arrow_left, arrow_top, arrow_width, arrow_height
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = dark_green
                arrow.line.color.rgb = dark_green

    else:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        apply_background(slide)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = slide_data["title"]
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = dark_green
        
        tf = content.text_frame
        tf.clear()
        for point in slide_data["content"]:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(40, 40, 40)

prs.save('Feature_Model_PoC.pptx')
print("Updated presentation with flowchart successfully.")
